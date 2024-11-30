import os
import json
import sys
import random
from pptx import Presentation
from win32com.client import Dispatch
from moviepy.editor import ImageClip, AudioFileClip

sys.path.append(os.getcwd())
from utils.enums import VIDEO_TYPE

content_path = os.path.join(os.getcwd(), 'create_content', 'created_content')
date_file_path = os.path.join(content_path, 'date_data.json')
content_file_path = os.path.join(content_path, 'content_data.json')
template_path = os.path.join(os.getcwd(), 'create_content', 'templates')

def create_image(creation_data):
    img_file_name = creation_data["post_id"] + "_img"
    content_details = creation_data["content_details"]
    content_keys = content_details.keys()
    img_content_path = os.path.join(content_path, 'images', img_file_name + '.pptx')
    match creation_data["post_type"]:
        case "word":
            img_template_path = os.path.join(template_path, creation_data["post_type"])
            prs = Presentation(os.path.join(img_template_path, "word_template_1.pptx"))

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Check special cases first (performance)
                            if (run.text == "grammarthree") & (content_details["word_type_id"].lower() != "verb"): run.text = "" 
                            elif (run.text == "word") &  (content_details["word_type_id"].lower() == "substantiv"): run.text = run.text.replace(run.text, content_details["article"] + " " + content_details["word"])
                            elif (run.text == "pronounciation"): run.text = run.text.replace(run.text, "[" + content_details["pronounciation"] + "]")
                            elif run.text in content_keys: 
                                    if "grammar" in run.text: run.text = content_details[run.text]["label"] + ": " + content_details[run.text]["value"]
                                    else: run.text = run.text.replace(run.text, content_details[run.text])
                                

            prs.save(img_content_path)
            application = Dispatch("PowerPoint.Application")
            pres = application.Presentations.Open(img_content_path)
            pres.Slides[0].Export(img_content_path.replace("pptx", "png"), "PNG")
            application.Quit()
            pres =  None
            application = None
             
    return creation_data["post_id"] + "_img"

def create_video(post_id):
    img_content_path = os.path.join(content_path, 'images', post_id + '_img.png')
    vid_content_path = os.path.join(content_path, 'videos', post_id + '_vid.mp4')
    music_content_path = os.path.join(template_path, 'music')

    music_tracks = os.listdir(music_content_path) 
    random_track = random.choice(music_tracks)
    audio = AudioFileClip(os.path.join(music_content_path, random_track)).set_duration(15)

    clip = ImageClip(img_content_path).set_duration(15).set_audio(audio)
    clip.write_videofile(vid_content_path, fps=24)
    return post_id + '_vid.mp4'

def create_content(upload_date, creation_data):

    data = {
        "post_id": creation_data["post_id"],
        "post_type": creation_data["post_type"],
        "upload_datetimeiso": upload_date + "T" + VIDEO_TYPE.WORD.UPLOAD_TIME,
        "tags": [],
        "vid_title": creation_data["content_details"]["word"] + " | Word of the Day",
        "post_title": "Word of the Day: " + creation_data["content_details"]["word"],
        "description": "TEST",
        "uploaded_all": False,
        "uploaded": {
            "youtube": False,
            "facebook": False,
            "twitter": False,
            "tiktok": False,
            "instagram": False
        }
    }

    data["image_id"] = create_image(creation_data)
    data["video_id"] = create_video(creation_data["post_id"])
    return data

def start_content_creation_process():

    days_json = {}
    content_json = {}

    if os.path.isfile(date_file_path):
        with open(date_file_path) as f:
            days_json = json.load(f)

    if os.path.isfile(content_file_path):
        with open(content_file_path) as f:
            content_json = json.load(f)
    
    for date in days_json:
        for content in days_json[date]["content"]:
            if content["content_created"] == False:
                if content_json.get("date") is None: content_json[date] = {}
                if content_json[date].get("content") is None: content_json[date]["content"] = []

                content_json[date]["content"].append(create_content(date, content))
                content["content_created"] == True


    json_object = json.dumps(days_json, indent=4)
    with open(date_file_path, "w") as outfile:
        outfile.write(json_object)
    
    json_object = json.dumps(content_json, indent=4)
    with open(content_file_path, "w") as outfile:
        outfile.write(json_object)